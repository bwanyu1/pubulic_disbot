import discord
import fitz  # PyMuPDF for PDF
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from groq import Groq
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer

# 環境変数を読み込む
load_dotenv()
DISCORD_TOKEN = os.getenv("DISCORD_TOKEN")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
EMBED_MODEL = os.getenv("EMBED_MODEL", "Qwen/qwen3-embedding-8b")
RERANK_MODEL = os.getenv("RERANK_MODEL", "Qwen/qwen3-4b-rerank")
PDF_CHUNK_SIZE = int(os.getenv("PDF_CHUNK_SIZE", 20))
PPTX_CHUNK_SIZE = int(os.getenv("PPTX_CHUNK_SIZE", 20))
TEXT_CHUNK_SIZE = int(os.getenv("TEXT_CHUNK_SIZE", 3000))
TOP_K = int(os.getenv("RAG_TOP_K", 5))

# インスタンス初期化
Groclient = Groq(api_key=GROQ_API_KEY)
embed_model = SentenceTransformer(EMBED_MODEL)
embed_dim = embed_model.get_sentence_embedding_dimension()
index = faiss.IndexFlatL2(embed_dim)
# FAISSインデックスに対応するメタデータリスト
metadata = []  # list of dict: {"thread_id": int, "text": str}

# Discord設定
intents = discord.Intents.default()
intents.message_content = True
client = discord.Client(intents=intents)
tree = discord.app_commands.CommandTree(client)

# --- ファイル抽出 & チャンク化関数 ---
def extract_and_chunk_pdf(file_bytes):
    chunks = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        pages = [page.get_text() for page in doc]
        for i in range(0, len(pages), PDF_CHUNK_SIZE):
            chunks.append("\n".join(pages[i:i+PDF_CHUNK_SIZE]))
    return chunks


def extract_and_chunk_pptx(file_bytes):
    chunks = []
    tmp = "temp.pptx"
    with open(tmp, "wb") as f:
        f.write(file_bytes)
    prs = Presentation(tmp)
    os.remove(tmp)
    slides = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        slides.append(text)
    for i in range(0, len(slides), PPTX_CHUNK_SIZE):
        chunks.append("\n".join(slides[i:i+PPTX_CHUNK_SIZE]))
    return chunks


def extract_and_chunk_docx(file_bytes):
    tmp = "temp.docx"
    with open(tmp, "wb") as f:
        f.write(file_bytes)
    doc = Document(tmp)
    os.remove(tmp)
    paras = [p.text for p in doc.paragraphs]
    text = "\n".join(paras)
    return [text[i:i+TEXT_CHUNK_SIZE] for i in range(0, len(text), TEXT_CHUNK_SIZE)]

def rerank_chunks(chunks):
    # Rerank用の埋め込み計算
    rerank_embeddings = embed_model.encode(chunks, convert_to_numpy=True)
    rerank_index = faiss.IndexFlatL2(embed_dim)
    rerank_index.add(rerank_embeddings.astype('float32'))
    
    # Rerank処理
    D, I = rerank_index.search(rerank_embeddings.astype('float32'), TOP_K)
    return [chunks[i] for i in I[0]]
def extract_and_chunk_txt(file_bytes):
    text = file_bytes.decode("utf-8")
    return [text[i:i+TEXT_CHUNK_SIZE] for i in range(0, len(text), TEXT_CHUNK_SIZE)]

# --- 要約 & FAISS登録 ---
async def process_and_store(chunks, thread_id):
    # 埋め込み計算
    embeddings = embed_model.encode(chunks, convert_to_numpy=True)
    # FAISSとmetadataに追加
    index.add(embeddings.astype('float32'))
    for text in chunks:
        metadata.append({"thread_id": thread_id, "text": text})

    # チャンクごとに要約
    summaries = []
    for idx, chunk in enumerate(chunks):
        prompt = (
            f"【チャンク {idx+1}/{len(chunks)}】\n"
            "以下の文章を要約し、日本語で出力してください。英語なら日本語に翻訳してください:\n\n" + chunk
        )
        resp = Groclient.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.3-70b-versatile",
        )
        summary = resp.choices[0].message.content.strip()
        summaries.append(summary)
        channel = client.get_channel(thread_id)
        await channel.send(f"**チャンク {idx+1} 要約**:\n{summary}")

    # 総合要約
    combined = "\n\n".join([f"チャンク {i+1}: {s}" for i, s in enumerate(summaries)])
    overall_prompt = (
        "以下は複数のチャンク要約です。これらを参考に全体の総合要約を作成し、日本語で出力してください。:\n\n"
        + combined
    )
    overall_resp = Groclient.chat.completions.create(
        messages=[{"role": "user", "content": overall_prompt}],
        model="llama-3.3-70b-versatile",
    )
    overall = overall_resp.choices[0].message.content.strip()
    channel = client.get_channel(thread_id)
    await channel.send(f"**💡 総合要約**:\n{overall}")

# --- メッセージ受信 ---
@client.event
async def on_message(message):
    if message.author.bot:
        return

    # 添付 or テキスト
    if message.attachments or message.content:
        # スレッド作成 or 取得
        if not isinstance(message.channel, discord.Thread):
            thread = await message.create_thread(name="要約・翻訳結果")
        else:
            thread = message.channel

        # テキスト単体
        if message.content and not message.attachments:
            chunks = [message.content]
        else:
            att = message.attachments[0]
            data = await att.read()
            name = att.filename.lower()
            if name.endswith('.pdf'):
                chunks = extract_and_chunk_pdf(data)
            elif name.endswith('.pptx'):
                chunks = extract_and_chunk_pptx(data)
            elif name.endswith('.docx'):
                chunks = extract_and_chunk_docx(data)
            elif name.endswith('.txt'):
                chunks = extract_and_chunk_txt(data)
            else:
                await thread.send("対応している形式は .txt, .pdf, .docx, .pptx です。")
                return

        # 処理 & 保存
        await process_and_store(chunks, thread.id)

# --- スラッシュコマンド: RAG QA ---
@tree.command(name="qa", description="アップロード済みドキュメントから回答を取得します。")
async def rag_qa(interaction: discord.Interaction, question: str):
    await interaction.response.defer()
    # 質問を埋め込み
    q_vec = embed_model.encode([question], convert_to_numpy=True).astype('float32')
    D, I = index.search(q_vec, TOP_K)
    # 上位チャンクを取得
    contexts = []
    for idx in I[0]:
        if idx < len(metadata):
            contexts.append(metadata[idx]['text'])

    # プロンプト作成
    prompt = (
        "以下のドキュメントの抜粋を参照して質問に回答してください。\n\n"
        + "\n\n".join(contexts)
        + f"\n\n質問: {question}"
    )
    resp = Groclient.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile",
    )
    answer = resp.choices[0].message.content.strip()
    await interaction.followup.send(f"**Q:** {question}\n**A:** {answer}")
@tree.command(name="image", description="画像をアップロードします。")
async def upload_image(interaction: discord.Interaction):
    if not interaction.message.attachments:
        await interaction.response.send_message("画像が添付されていません。", ephemeral=True)
        return

    att = interaction.message.attachments[0]
    data = await att.read()
    image_url = att.url

    # 画像をGroqに送信
    resp = Groclient.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "この画像は何ですか？"},
                    {"type": "image_url", "image_url": {"url": image_url}}
                ]
            }
        ],
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        temperature=1,
        max_completion_tokens=1024,
        top_p=1,
        stream=False,
        stop=None,
    )
    answer = resp.choices[0].message.content.strip()
    await interaction.response.send_message(f"**回答:** {answer}")
@client.event
async def on_ready():
    await tree.sync()
    print(f"Logged in as {client.user}")

# Bot起動
client.run(DISCORD_TOKEN)
